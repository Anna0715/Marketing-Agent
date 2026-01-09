"""文档加载器 - 支持多种格式"""
import os
from pathlib import Path
from typing import List, Dict, Any, Optional
from abc import ABC, abstractmethod
import PyPDF2 
from pptx import Presentation  
from docx import Document 
import requests  
from datetime import datetime

from src.models.schemas import DocumentType, DocumentMetadata
from src.utils.config import settings


class BaseDocumentLoader(ABC):
    """文档加载器基类"""
    
    @abstractmethod
    def load(self, file_path: str) -> List[Dict[str, Any]]:
        """加载文档并返回文本块列表"""
        pass
    
    def extract_metadata(self, file_path: str, doc_type: DocumentType) -> DocumentMetadata:
        """提取文档元数据"""
        path_obj = Path(file_path)
        return DocumentMetadata(
            doc_id=f"{doc_type.value}_{path_obj.stem}_{int(datetime.now().timestamp())}",
            title=path_obj.stem,
            doc_type=doc_type,
            upload_time=datetime.now(),
            file_path=str(path_obj.absolute()),
            tags=[],
            metadata={"file_size": path_obj.stat().st_size if path_obj.exists() else 0}
        )


class PDFLoader(BaseDocumentLoader):
    """PDF文档加载器"""
    
    def load(self, file_path: str) -> List[Dict[str, Any]]:
        """加载PDF文件"""
        chunks = []
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                full_text = ""
                
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    text = page.extract_text()
                    full_text += f"\n\n--- 第 {page_num} 页 ---\n\n{text}"
                
                # 简单的分块策略
                chunk_size = settings.chunk_size
                for i in range(0, len(full_text), chunk_size - settings.chunk_overlap):
                    chunk = full_text[i:i + chunk_size]
                    if chunk.strip():
                        chunks.append({
                            "content": chunk.strip(),
                            "metadata": {
                                "source": file_path,
                                "page": page_num,
                                "chunk_index": len(chunks)
                            }
                        })
        except Exception as e:
            print(f"加载PDF文件出错: {e}")
            raise
        
        return chunks


class PPTXLoader(BaseDocumentLoader):
    """PPTX文档加载器"""
    
    def load(self, file_path: str) -> List[Dict[str, Any]]:
        """加载PPTX文件"""
        chunks = []
        try:
            prs = Presentation(file_path)
            full_text = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"\n\n--- 幻灯片 {slide_num} ---\n\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                
                full_text += slide_text
            
            # 按幻灯片分块
            chunk_size = settings.chunk_size
            for i in range(0, len(full_text), chunk_size - settings.chunk_overlap):
                chunk = full_text[i:i + chunk_size]
                if chunk.strip():
                    chunks.append({
                        "content": chunk.strip(),
                        "metadata": {
                            "source": file_path,
                            "chunk_index": len(chunks)
                        }
                    })
        except Exception as e:
            print(f"加载PPTX文件出错: {e}")
            raise
        
        return chunks


class DOCXLoader(BaseDocumentLoader):
    """DOCX文档加载器"""
    
    def load(self, file_path: str) -> List[Dict[str, Any]]:
        """加载DOCX文件"""
        chunks = []
        try:
            doc = Document(file_path)
            full_text = "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
            
            chunk_size = settings.chunk_size
            for i in range(0, len(full_text), chunk_size - settings.chunk_overlap):
                chunk = full_text[i:i + chunk_size]
                if chunk.strip():
                    chunks.append({
                        "content": chunk.strip(),
                        "metadata": {
                            "source": file_path,
                            "chunk_index": len(chunks)
                        }
                    })
        except Exception as e:
            print(f"加载DOCX文件出错: {e}")
            raise
        
        return chunks


class TextLoader(BaseDocumentLoader):
    """文本文件加载器"""
    
    def load(self, file_path: str) -> List[Dict[str, Any]]:
        """加载文本文件"""
        chunks = []
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            chunk_size = settings.chunk_size
            for i in range(0, len(content), chunk_size - settings.chunk_overlap):
                chunk = content[i:i + chunk_size]
                if chunk.strip():
                    chunks.append({
                        "content": chunk.strip(),
                        "metadata": {
                            "source": file_path,
                            "chunk_index": len(chunks)
                        }
                    })
        except Exception as e:
            print(f"加载文本文件出错: {e}")
            raise
        
        return chunks


class FeishuLoader(BaseDocumentLoader):
    """飞书文档加载器"""
    
    def __init__(self):
        self.base_url = "https://open.feishu.cn"
        self.app_id = settings.feishu_app_id
        self.app_secret = settings.feishu_app_secret
        self.access_token = None
        self._get_access_token()
    
    def _get_access_token(self):
        """获取飞书访问令牌"""
        if not self.app_id or not self.app_secret:
            print("警告: 未配置飞书App ID和Secret，无法加载飞书文档")
            return
        
        try:
            url = f"{self.base_url}/open-apis/auth/v3/tenant_access_token/internal"
            headers = {"Content-Type": "application/json; charset=utf-8"}
            data = {
                "app_id": self.app_id,
                "app_secret": self.app_secret
            }
            response = requests.post(url, json=data, headers=headers, timeout=10)
            if response.status_code == 200:
                result = response.json()
                if result.get("code") == 0:
                    self.access_token = result.get("tenant_access_token")
        except Exception as e:
            print(f"获取飞书访问令牌失败: {e}")
    
    def load(self, file_path: str) -> List[Dict[str, Any]]:
        """加载飞书文档
        file_path格式: feishu://doc_token 或 feishu://document_id
        """
        chunks = []
        
        if not self.access_token:
            print("无法访问飞书文档：缺少访问令牌")
            return chunks
        
        try:
            # 提取文档token或ID
            doc_token = file_path.replace("feishu://", "")
            
            # 调用飞书API获取文档内容
            url = f"{self.base_url}/open-apis/drive/v1/files/{doc_token}/content"
            headers = {
                "Authorization": f"Bearer {self.access_token}",
                "Content-Type": "application/json"
            }
            
            response = requests.get(url, headers=headers, timeout=30)
            if response.status_code == 200:
                result = response.json()
                if result.get("code") == 0:
                    content = result.get("data", {}).get("content", "")
                    
                    # 处理内容分块
                    chunk_size = settings.chunk_size
                    for i in range(0, len(content), chunk_size - settings.chunk_overlap):
                        chunk = content[i:i + chunk_size]
                        if chunk.strip():
                            chunks.append({
                                "content": chunk.strip(),
                                "metadata": {
                                    "source": file_path,
                                    "doc_token": doc_token,
                                    "chunk_index": len(chunks)
                                }
                            })
        except Exception as e:
            print(f"加载飞书文档出错: {e}")
        
        return chunks


class DocumentLoaderFactory:
    """文档加载器工厂"""
    
    _loaders = {
        DocumentType.PDF: PDFLoader(),
        DocumentType.PPTX: PPTXLoader(),
        DocumentType.DOCX: DOCXLoader(),
        DocumentType.TXT: TextLoader(),
        DocumentType.MD: TextLoader(),
        DocumentType.FEISHU: FeishuLoader(),
    }
    
    @classmethod
    def get_loader(cls, doc_type: DocumentType) -> BaseDocumentLoader:
        """根据文档类型获取对应的加载器"""
        loader = cls._loaders.get(doc_type)
        if not loader:
            raise ValueError(f"不支持的文档类型: {doc_type}")
        return loader
    
    @classmethod
    def detect_doc_type(cls, file_path: str) -> DocumentType:
        """检测文档类型"""
        path_obj = Path(file_path)
        extension = path_obj.suffix.lower()
        
        if file_path.startswith("feishu://"):
            return DocumentType.FEISHU
        
        type_map = {
            ".pdf": DocumentType.PDF,
            ".pptx": DocumentType.PPTX,
            ".docx": DocumentType.DOCX,
            ".txt": DocumentType.TXT,
            ".md": DocumentType.MD,
        }
        
        return type_map.get(extension, DocumentType.TXT)
